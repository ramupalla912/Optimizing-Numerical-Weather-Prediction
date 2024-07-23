# # importing libraries 
# import pandas as pd 
# import sklearn
# import numpy as np 
# import sklearn as sk 
# from sklearn.linear_model import LinearRegression 
# import matplotlib.pyplot as plt 

# # read the cleaned data 
# data = pd.read_csv("austin_final.csv") 

# # the features or the 'x' values of the data 
# # these columns are used to train the model 
# # the last column, i.e, precipitation column 
# # will serve as the label 
# X = data.drop(['PrecipitationSumInches'], axis = 1) 

# # the output or the label. 
# Y = data['PrecipitationSumInches'] 
# # reshaping it into a 2-D vector 

# Y = Y.values.reshape(-1, 1) 

# # consider a random day in the dataset 
# # we shall plot a graph and observe this 
# # day 

# day_index = 798
# days = [i for i in range(Y.size)] 

# # initialize a linear regression classifier 

# clf = LinearRegression() 
# # train the classifier with our 
# # input data. 

# clf.fit(X, Y) 

# # give a sample input to test our model 
# # this is a 2-D vector that contains values 
# # for each column in the dataset. 

# inp = np.array([[74], [60], [45], [67], [49], [43], [33], [45], 
# 				[57], [29.68], [10], [7], [2], [0], [20], [4], [31]]) 
# inp = inp.reshape(1, -1) 

# # print the output. 

# print('The precipitation in inches for the input is:', clf.predict(inp)) 

# # plot a graph of the precipitation levels 
# # versus the total number of days. 
# # one day, which is in red, is 
# # tracked here. It has a precipitation 
# # of approx. 2 inches. 

# print("the precipitation trend graph: ") 
# plt.scatter(days, Y, color = 'g') 
# plt.scatter(days[day_index], Y[day_index], color ='r') 
# plt.title("Precipitation level") 
# plt.xlabel("Days") 
# plt.ylabel("Precipitation in inches") 


# plt.show() 
# x_vis = X.filter(['TempAvgF', 'DewPointAvgF', 'HumidityAvgPercent', 
# 				'SeaLevelPressureAvgInches', 'VisibilityAvgMiles', 
# 				'WindAvgMPH'], axis = 1) 

# # plot a graph with a few features (x values) 
# # against the precipitation or rainfall to observe 
# # the trends 



# print("Precipitation vs selected attributes graph: ") 

# for i in range(x_vis.columns.size): 
# 	plt.subplot(3, 2, i + 1) 
# 	plt.scatter(days, x_vis[x_vis.columns.values[i][:100]], 
# 											color = 'g') 

# 	plt.scatter(days[day_index], 
# 				x_vis[x_vis.columns.values[i]][day_index], 
# 				color ='r') 

# 	plt.title(x_vis.columns.values[i]) 

# plt.show() 







# # import pandas as pd 
# # import sklearn
# # import numpy as np 
# # import sklearn as sk 
# # from sklearn.linear_model import LinearRegression 
# # import matplotlib.pyplot as plt 

# # data = pd.read_csv("austin_final.csv") 
# # X = data.drop(['PrecipitationSumInches'], axis = 1) 

# # Y = data['PrecipitationSumInches'] 
# # Y = Y.values.reshape(-1, 1) 

# # day_index = 798
# # days = [i for i in range(Y.size)] 
# # clf = LinearRegression() 
# # clf.fit(X, Y) 
# # inp = np.array([[74], [60], [45], [67], [49], [43], [33], [45], 
# # 				[57], [29.68], [10], [7], [2], [0], [20], [4], [31]]) 
# # inp = inp.reshape(1, -1) 
# # print('The precipitation in inches for the input is:', clf.predict(inp)) 
# # print("the precipitation trend graph: ") 
# # plt.scatter(days, Y, color = 'g') 
# # plt.scatter(days[day_index], Y[day_index], color ='r') 
# # plt.title("Precipitation level") 
# # plt.xlabel("Days") 
# # plt.ylabel("Precipitation in inches") 


# # plt.show() 
# # x_vis = X.filter(['TempAvgF', 'DewPointAvgF', 'HumidityAvgPercent', 
# # 				'SeaLevelPressureAvgInches', 'VisibilityAvgMiles', 
# # 				'WindAvgMPH'], axis = 1) 

# # print("Precipitation vs selected attributes graph: ") 

# # for i in range(x_vis.columns.size): 
# # 	plt.subplot(3, 2, i + 1) 
# # 	plt.scatter(days, x_vis[x_vis.columns.values[i][:100]], 
# # 											color = 'g') 

# # 	plt.scatter(days[day_index], 
# # 				x_vis[x_vis.columns.values[i]][day_index], 
# # 				color ='r') 

# # 	plt.title(x_vis.columns.values[i]) 

# # plt.show() 



from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Optimizing Numerical Weather Prediction Model Performance Using Machine Learning Techniques"
subtitle.text = "Team Members:\nS. Pavan (21R01A05P7)\nK. Hemanth Kumar (21R01A05M5)\nP. Ramu (22R05A0520)\nS. Dhanunjay (22R05A0521)"

# Slide 2: Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Introduction"
content = slide.placeholders[1]
content.text = "This project aims to enhance the performance of Numerical Weather Prediction (NWP) models using machine learning techniques."

# Slide 3: Data and Methods
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Data and Methods"
content = slide.placeholders[1]
content.text = ("The data for this project is sourced from the 'austin_final.csv' file.\n"
                "We use a Linear Regression model to predict precipitation levels based on various meteorological parameters.\n"
                "The code is implemented using Python with libraries such as pandas, numpy, and scikit-learn.")

# Slide 4: Code Snippet
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Code Snippet"
content = slide.placeholders[1]
content.text = ("import pandas as pd\n"
                "import numpy as np\n"
                "from sklearn.linear_model import LinearRegression\n"
                "import matplotlib.pyplot as plt\n\n"
                "data = pd.read_csv('austin_final.csv')\n"
                "X = data.drop(['PrecipitationSumInches'], axis = 1)\n"
                "Y = data['PrecipitationSumInches'].values.reshape(-1, 1)\n"
                "clf = LinearRegression().fit(X, Y)\n"
                "inp = np.array([[74], [60], [45], [67], [49], [43], [33], [45],\n"
                "                [57], [29.68], [10], [7], [2], [0], [20], [4], [31]]).reshape(1, -1)\n"
                "print('The precipitation in inches for the input is:', clf.predict(inp))")

# Slide 5: Code Snippet (contd.)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Code Snippet (contd.)"
content = slide.placeholders[1]
content.text = ("print('The precipitation trend graph: ')\n"
                "plt.scatter(days, Y, color = 'g')\n"
                "plt.scatter(days[day_index], Y[day_index], color ='r')\n"
                "plt.title('Precipitation level')\n"
                "plt.xlabel('Days')\n"
                "plt.ylabel('Precipitation in inches')\n"
                "plt.show()\n"
                "x_vis = X.filter(['TempAvgF', 'DewPointAvgF', 'HumidityAvgPercent',\n"
                "                  'SeaLevelPressureAvgInches', 'VisibilityAvgMiles',\n"
                "                  'WindAvgMPH'], axis = 1)\n"
                "print('Precipitation vs selected attributes graph: ')\n"
                "for i in range(x_vis.columns.size):\n"
                "    plt.subplot(3, 2, i + 1)\n"
                "    plt.scatter(days, x_vis[x_vis.columns.values[i][:100]], color = 'g')\n"
                "    plt.scatter(days[day_index], x_vis[x_vis.columns.values[i]][day_index], color ='r')\n"
                "    plt.title(x_vis.columns.values[i])\n"
                "plt.show()")

# Slide 6: Results - Precipitation Prediction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Results - Precipitation Prediction"
content = slide.placeholders[1]
content.text = "The model predicts the precipitation in inches for the given input."

# Slide 7: Results - Precipitation Trend Graph
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Results - Precipitation Trend Graph"
content = slide.placeholders[0]
content.text = "Insert precipitation trend graph here."

# Slide 8: Results - Attribute Analysis
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Results - Attribute Analysis"
content = slide.placeholders[1]
content.text = ("The following graphs show the relationship between precipitation and selected attributes:\n"
                "- TempAvgF\n"
                "- DewPointAvgF\n"
                "- HumidityAvgPercent\n"
                "- SeaLevelPressureAvgInches\n"
                "- VisibilityAvgMiles\n"
                "- WindAvgMPH")

# Slide 9: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1]
content.text = ("The Linear Regression model shows a relationship between meteorological parameters and precipitation levels.\n"
                "Further optimization and testing with different machine learning models can enhance prediction accuracy.")

# Slide 10: Questions
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Questions"
content = slide.placeholders[1]
content.text = "Thank you! Any questions?"

# Save the presentation
pptx_file = "/mnt/data/Numerical_Weather_Prediction_Presentation.pptx"
prs.save(pptx_file)

pptx_file
